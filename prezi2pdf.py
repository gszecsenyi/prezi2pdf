import requests
import re, os, json, sys
from io import BytesIO
from img2pdf import convert
import yt_dlp
import argparse
from PIL import Image
from pptx import Presentation


# ArgParse
parser = argparse.ArgumentParser(description='Download Prezi Presentations and Videos')
parser.add_argument('--url','-u',dest='url', action='store', help='Prezi URL', required=True)
parser.add_argument('--download-json','-j',dest='download_json', action='store_true', help='Download JSON file', required=False)
parser.add_argument('--output-type','-t',dest='output_type', choices=['pdf','pptx'], default='pdf', help='Choose output format for presentations', required=False)

args = parser.parse_args()


def download_video(id):
    
    url = f"https://prezi.com/api/v5/presentation-content/{id}/"
    data = requests.get(url).json()
    try:
        os.mkdir(f"./videos")
    except FileExistsError:
        pass
    print(f"Downloading {data['meta']['title']}")
    video_url = data['meta']['video_signed_url_with_title']
    ydl = yt_dlp.YoutubeDL({'outtmpl': f'./videos/{id}.%(ext)s',
                            'merge_output_format': 'mp4', 'ignoreerrors': 'True',
                            'writethumbnail': 'True', 'retries': 100,
                            'add-header': 'user-agent:Mozilla/5.0',})
    info_dict = ydl.extract_info(video_url, download=True)
    if args.download_json:
        with open(f"./videos/{id}.json", 'w') as outfile:
            outfile.writelines(json.dumps(data, indent=4))


def save_presentation_pdf(content, id):
    with open(f'presentations/{id}.pdf', 'wb') as pdf:
        pdf.write(convert(content))


def save_presentation_ppt(content, id):
    first_image = Image.open(BytesIO(content[0]))
    width_px, height_px = first_image.size
    first_image.close()
    emu_per_px = 914400 / 96
    slide_width = int(width_px * emu_per_px)
    slide_height = int(height_px * emu_per_px)
    prs = Presentation()
    prs.slide_width = slide_width
    prs.slide_height = slide_height
    blank_layout = prs.slide_layouts[6]
    for page_data in content:
        slide = prs.slides.add_slide(blank_layout)
        image_stream = BytesIO(page_data)
        slide.shapes.add_picture(
            image_stream,
            0,
            0,
            width=slide_width,
            height=slide_height,
        )
    prs.save(f'presentations/{id}.pptx')


def download_presentation(id):
    url = f"https://prezi.com/api/v2/storyboard/{id}/"
    print("Requesting data from API, please wait up to 60 seconds...")
    data = requests.get(url, timeout=60).json()
    try:
        os.mkdir(f"./presentations")
    except FileExistsError:
        pass
    content = []
    i = 0
    total = len(data['steps'])
    for frame in data['steps']:
        r = requests.get(frame['images'][0]['url'])
        content.append(r.content)
        print(f"Downloading slide {i+1}/{total}")
        i += 1
    if args.output_type == 'pdf':
        save_presentation_pdf(content, id)
    else:
        save_presentation_ppt(content, id)
    if args.download_json:
        with open(f"./presentations/{id}.json", 'w') as outfile:
            outfile.writelines(json.dumps(data, indent=4))

match = re.search(r"prezi\.com/(?:p/(?:edit/)?)?([0-9a-zA-Z-]{12})", args.url)
if not match:
    print("Please provide a valid Prezi URL that contains the 12-character presentation ID.")
    sys.exit(1)

id = match.group(1)

if "prezi.com/v/" in args.url:
    download_video(id)

elif "prezi.com/i/" in args.url:
    print("Prezi design not supported yet")

elif "prezi.com/" in args.url:
    download_presentation(id)
else:
    print("Please provide a valid prezi URL")
